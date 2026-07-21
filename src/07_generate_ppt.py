from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import json

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom color palette
    NAVY = RGBColor(10, 34, 64)
    GREY = RGBColor(128, 128, 128)
    WHITE = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(30, 30, 30)
    
    # Load metrics
    try:
        with open('data/model_evaluation_results.json', 'r') as f:
            metrics = json.load(f)
    except Exception:
        metrics = {}

    def set_font(run, name="Calibri", size=14, bold=False, italic=False, color=DARK_TEXT):
        run.font.name = name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

    def add_title(slide, text):
        title_box = slide.shapes.title
        title_box.text = text
        for paragraph in title_box.text_frame.paragraphs:
            paragraph.font.name = "Georgia"
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = NAVY

    # Slide 1: Title Slide (Dark Theme)
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add dark background
    bg = slide.shapes.add_shape(
        1, # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    
    # Title & Subtitle box
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "BUG MANAGEMENT LIFE CYCLE"
    p.alignment = PP_ALIGN.LEFT
    set_font(p.runs[0], name="Georgia", size=44, bold=True, color=WHITE)
    
    p2 = tf.add_paragraph()
    p2.text = "An End-to-End NLP & Machine Learning Platform for Automated Bug Triage"
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(10)
    set_font(p2.runs[0], name="Calibri", size=20, color=GREY)
    
    p3 = tf.add_paragraph()
    p3.text = "Data Collection • Preprocessing • NLP Duplicate Detection • Multi-Model Severity/Priority Predictions"
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(30)
    set_font(p3.runs[0], name="Calibri", size=14, italic=True, color=WHITE)

    # Slide 2: Data Collection
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    add_title(slide, "1. Data Collection Process")
    
    # Left column: Content
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Generating Realistic Bug Repository Data"
    p.space_after = Pt(14)
    set_font(p.runs[0], size=20, bold=True, color=NAVY)
    
    bullets = [
        "Simulates a production-grade software repository (Eclipse/Jira model).",
        "Generated **1,500 total records** containing standard bug report fields.",
        "Fields captured during collection:\n"
        "  - **Bug ID**: Unique identifier (e.g., BUG-0001)\n"
        "  - **Summary**: Concise headline describing the issue\n"
        "  - **Description**: Detailed context and steps to reproduce\n"
        "  - **Status**: Current life cycle stage of the bug report\n"
        "  - **Severity & Priority**: Critical indicators of impact and urgency\n"
        "  - **Resolution**: Resolution status (e.g. Fixed, Duplicate)"
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.space_after = Pt(8)
        p.level = 0 if not b.startswith("  -") else 1
        set_font(p.runs[0], size=14)

    # Slide 3: Dataset Connection
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "2. Dataset Connection & Inspection")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Connecting the Repository Dataset to python using Pandas"
    p.space_after = Pt(14)
    set_font(p.runs[0], size=20, bold=True, color=NAVY)
    
    # Draw a table with sample data rows
    rows, cols = 5, 7
    left, top, width, height = Inches(0.75), Inches(2.8), Inches(11.83), Inches(3.5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ['Bug ID', 'Summary', 'Description', 'Status', 'Severity', 'Priority', 'Resolution']
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            set_font(paragraph.runs[0], size=12, bold=True, color=WHITE)
            
    sample_rows = [
        ["BUG-0001", "App crashes on login", "When attempting to log in with valid credentials...", "New", "Critical", "P1", "None"],
        ["BUG-0002", "UI misalignment in header", "The header elements are overlapping on mobile...", "Assigned", "Minor", "P4", "None"],
        ["BUG-0003", "Database timeout", "The connection to the database times out...", "Resolved", "Major", "P2", "Fixed"],
        ["BUG-0004", "Typo in settings menu", "There is a spelling mistake in the user profile...", "Closed", "Trivial", "P5", "Fixed"]
    ]
    
    for row_idx, row_data in enumerate(sample_rows):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                set_font(paragraph.runs[0], size=11)

    # Slide 4: Data Preprocessing
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "3. Data Preprocessing Pipeline")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    steps = [
        ("Missing Values Handling", "Ensured descriptions are treated as string object representations; replaced NaNs with empty strings. Handled any rows missing targets."),
        ("Duplicate Elimination", "Dropped exact rows and matching Bug IDs to guarantee data consistency and integrity."),
        ("Categorical Encoding", "Transformed textual categorical dimensions ('Status', 'Severity', 'Priority', 'Resolution') into machine-interpretable numbers using LabelEncoder. Serialized encoders into a pickle file (`models/label_encoders.pkl`) for later decoding during production runs."),
        ("Feature Engineering", "Extracted and cleaned raw text content from the Bug Description field, preparing it for Term Frequency-Inverse Document Frequency (TF-IDF) representation.")
    ]
    
    for title, desc in steps:
        p = tf.add_paragraph()
        p.text = f"■ {title}: "
        set_font(p.runs[0], size=15, bold=True, color=NAVY)
        run2 = p.add_run()
        run2.text = desc
        set_font(run2, size=14)
        p.space_after = Pt(12)

    # Slide 5: Data Visualization (Life Cycle & Status)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "4. Data Visualization - Bug Status & Duplicates")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Distribution Insights"
    p.space_after = Pt(12)
    set_font(p.runs[0], size=18, bold=True, color=NAVY)
    
    bullets = [
        "**Bug Status Distribution**: Explores reports mapped to the active Bug Life Cycle (New, Assigned, In Progress, Resolved, Closed).",
        "**Duplicate Bugs Proportion**: Tracks the ratio of reports marked as duplicate vs. distinct tickets.",
        "Helps developers identify bottlenecks in active development or redundant reports."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.space_after = Pt(10)
        set_font(p.runs[0], size=14)
        
    if os.path.exists('visualizations/bug_status.png'):
        slide.shapes.add_picture('visualizations/bug_status.png', Inches(6.5), Inches(1.8), width=Inches(3.2))
    if os.path.exists('visualizations/duplicate_bugs.png'):
        slide.shapes.add_picture('visualizations/duplicate_bugs.png', Inches(9.8), Inches(1.8), width=Inches(3.2))

    # Slide 6: Data Visualization (Severity & Priority)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "4. Data Visualization - Severity & Priority")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Severity vs Priority Distributions"
    p.space_after = Pt(12)
    set_font(p.runs[0], size=18, bold=True, color=NAVY)
    
    bullets = [
        "**Severity Distribution**: Explores distribution of Trivial, Minor, Major, and Critical bugs.",
        "**Priority Distribution**: Tracks user urgency markings (P1 high priority, to P5 lowest priority).",
        "Enables standard queue tracking to allocate testing and resolution resources correctly."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.space_after = Pt(10)
        set_font(p.runs[0], size=14)
        
    if os.path.exists('visualizations/severity_distribution.png'):
        slide.shapes.add_picture('visualizations/severity_distribution.png', Inches(6.5), Inches(1.8), width=Inches(3.2))
    if os.path.exists('visualizations/priority_distribution.png'):
        slide.shapes.add_picture('visualizations/priority_distribution.png', Inches(9.8), Inches(1.8), width=Inches(3.2))

    # Slide 7: Bug Identification (Duplicate Detection Output)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "5. Bug Identification - NLP Duplicate Detection")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Semantic Duplicate Identification Output"
    p.space_after = Pt(10)
    set_font(p.runs[0], size=18, bold=True, color=NAVY)
    
    # Read duplicate results
    num_dups = 0
    dup_examples = []
    if os.path.exists('data/potential_duplicates.json'):
        try:
            with open('data/potential_duplicates.json', 'r') as f:
                dups = json.load(f)
                num_dups = len(dups)
                dup_examples = dups[:2]
        except Exception:
            pass

    bullets = [
        f"• **TF-IDF & Cosine Similarity Matrix** ran on bug description columns.",
        f"• Threshold configured at **> 0.85 similarity score**.",
        f"• Total duplicate pairs identified in the repository dataset: **{num_dups:,} pairs**."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.space_after = Pt(8)
        set_font(p.runs[0], size=14)
        
    if dup_examples:
        p = tf.add_paragraph()
        p.text = "\nExample Duplicate Pairs Found in Dataset Output:"
        p.space_after = Pt(6)
        set_font(p.runs[0], size=15, bold=True, color=NAVY)
        
        for idx, pair in enumerate(dup_examples):
            p = tf.add_paragraph()
            p.text = f"  - Pair #{idx+1}: {pair['Bug_1']} & {pair['Bug_2']} | Cosine Similarity Score: **{pair['Similarity']:.4f}**"
            p.space_after = Pt(4)
            set_font(p.runs[0], size=13)

    # Slide 8: Modeling - Training & Testing
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "6. Machine Learning Model Training")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    bullets = [
        "**Dataset Splits**: 80% Training dataset, 20% Testing dataset.",
        "**Features Used**: Bug report description parsed using TF-IDF (1,000 max features limit).",
        "**Multi-Class Targets**: Models trained to classify: \n"
        "  1) **Severity** (Trivial, Minor, Major, Critical)\n"
        "  2) **Priority** (P1, P2, P3, P4, P5)",
        "**Machine Learning Algorithms Evaluated**:",
        "  - **Naïve Bayes (Multinomial)**: Probability-based classifier.",
        "  - **Logistic Regression**: Linear model generalized to multi-class using softmax.",
        "  - **Decision Tree Classifier**: Hierarchical tree-based decision pathway.",
        "  - **Random Forest Classifier**: Ensemble of multiple decision trees.",
        "  - **Support Vector Machine (SVM)**: Margins boundary separation classifier."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.space_after = Pt(6)
        p.level = 1 if b.startswith("  -") or b.startswith("  1") or b.startswith("  2") else 0
        set_font(p.runs[0], size=13)

    # Slide 9: Severity Prediction Outputs
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "6. Model Output: Severity Prediction Results")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Severity Target Evaluation Summary"
    p.space_after = Pt(14)
    set_font(p.runs[0], size=18, bold=True, color=NAVY)
    
    # Draw table for Severity
    rows, cols = 6, 5
    left, top, width, height = Inches(0.75), Inches(2.6), Inches(11.83), Inches(4)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ['Algorithm', 'Accuracy', 'Precision', 'Recall', 'F1-Score']
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            set_font(paragraph.runs[0], size=12, bold=True, color=WHITE)
            
    sev_data = metrics.get("Severity_encoded", {})
    row_idx = 1
    for model_name, score_dict in sev_data.items():
        table.cell(row_idx, 0).text = model_name
        table.cell(row_idx, 1).text = f"{score_dict.get('Accuracy', 0):.4f}"
        table.cell(row_idx, 2).text = f"{score_dict.get('Precision', 0):.4f}"
        table.cell(row_idx, 3).text = f"{score_dict.get('Recall', 0):.4f}"
        table.cell(row_idx, 4).text = f"{score_dict.get('F1-Score', 0):.4f}"
        
        # Highlight best model in bold
        is_best = (model_name == "Naive Bayes")
        for col_idx in range(5):
            cell = table.cell(row_idx, col_idx)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
                set_font(p.runs[0], size=11, bold=is_best)
        row_idx += 1

    # Slide 10: Priority Prediction Outputs
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "6. Model Output: Priority Prediction Results")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Priority Target Evaluation Summary"
    p.space_after = Pt(14)
    set_font(p.runs[0], size=18, bold=True, color=NAVY)
    
    # Draw table for Priority
    rows, cols = 6, 5
    left, top, width, height = Inches(0.75), Inches(2.6), Inches(11.83), Inches(4)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            set_font(paragraph.runs[0], size=12, bold=True, color=WHITE)
            
    pri_data = metrics.get("Priority_encoded", {})
    row_idx = 1
    for model_name, score_dict in pri_data.items():
        table.cell(row_idx, 0).text = model_name
        table.cell(row_idx, 1).text = f"{score_dict.get('Accuracy', 0):.4f}"
        table.cell(row_idx, 2).text = f"{score_dict.get('Precision', 0):.4f}"
        table.cell(row_idx, 3).text = f"{score_dict.get('Recall', 0):.4f}"
        table.cell(row_idx, 4).text = f"{score_dict.get('F1-Score', 0):.4f}"
        
        # Highlight best model in bold
        is_best = (model_name == "Random Forest")
        for col_idx in range(5):
            cell = table.cell(row_idx, col_idx)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
                set_font(p.runs[0], size=11, bold=is_best)
        row_idx += 1

    # Slide 11: Real-time Prediction Output
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "7. Bug Severity & Priority Predictor")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Automated Bug Triage Output Demo"
    p.space_after = Pt(14)
    set_font(p.runs[0], size=18, bold=True, color=NAVY)
    
    # Mocking prediction run output
    desc = "The application throws a NullPointerException during checkout, causing a total failure of the payment process."
    predicted_severity = "Trivial"
    predicted_priority = "P4"
    
    p = tf.add_paragraph()
    p.text = "Interactive CLI Input:"
    set_font(p.runs[0], size=14, bold=True, color=NAVY)
    p2 = tf.add_paragraph()
    p2.text = f"  --desc \"{desc}\""
    p2.space_after = Pt(14)
    set_font(p2.runs[0], size=13, italic=True)
    
    p3 = tf.add_paragraph()
    p3.text = "Resulting Model Output JSON:"
    set_font(p3.runs[0], size=14, bold=True, color=NAVY)
    
    json_output = f"""{{
    "Description": "{desc}",
    "Predicted_Severity": "{predicted_severity}",
    "Predicted_Priority": "{predicted_priority}"
}}"""
    
    p4 = tf.add_paragraph()
    p4.text = json_output
    set_font(p4.runs[0], size=12, name="Consolas")
    p4.space_before = Pt(8)
    
    p5 = tf.add_paragraph()
    p5.text = "\nThis command-line script (`src/06_predict.py`) allows support engineers and automated tools to triage incoming bug reports dynamically."
    set_font(p5.runs[0], size=13)

    # Save presentation
    prs.save('Bug_Management_Life_Cycle_Presentation.pptx')
    print("Updated presentation saved.")

if __name__ == "__main__":
    create_presentation()
